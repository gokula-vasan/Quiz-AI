import fitz  # PyMuPDF
import re

def clean_text(text: str) -> str:
    """
    Cleans extracted text by removing control characters, resolving duplicate spacing, 
    and standardizing paragraph breaks.
    """
    # 1. Remove weird control bytes/characters (keep standard tabs/newlines)
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    # 2. Convert multiple consecutive spaces to a single space
    text = re.sub(r'[ \t]+', ' ', text)
    # 3. Clean duplicate empty lines
    text = re.sub(r'\n\s*\n', '\n\n', text)
    return text.strip()

def parse_pdf(file_bytes: bytes) -> str:
    """
    Opens a PDF file from raw bytes using PyMuPDF and extracts its text contents.
    """
    text_content = ""
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            text_content += page.get_text() + "\n"
        doc.close()
    except Exception as e:
        raise ValueError(f"Error parsing PDF file: {str(e)}")
    return clean_text(text_content)

def parse_pptx(file_bytes: bytes) -> str:
    """
    Opens a PPTX file from raw bytes using python-pptx and extracts text from shapes on all slides.
    """
    text_content = ""
    try:
        import io
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                # Extract text from standard text frames (titles, body content)
                if shape.has_text_frame:
                    if shape.text_frame.text.strip():
                        text_content += shape.text_frame.text + "\n"
                # Extract text from slide tables if present
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text_content += cell.text + "\n"
    except Exception as e:
        raise ValueError(f"Error parsing PowerPoint file: {str(e)}")
    return clean_text(text_content)

def parse_docx(file_bytes: bytes) -> str:
    """
    Opens a DOCX file from raw bytes using python-docx and extracts text from paragraphs and tables.
    """
    text_content = ""
    try:
        import io
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            if para.text.strip():
                text_content += para.text + "\n"
        # Also extract text from tables inside the Word document
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_content += cell.text + "\n"
    except Exception as e:
        raise ValueError(f"Error parsing Word document: {str(e)}")
    return clean_text(text_content)

def parse_txt(file_bytes: bytes) -> str:
    """
    Decodes plain text notes from raw bytes, supporting UTF-8 and Latin-1 fallbacks.
    """
    try:
        text_content = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text_content = file_bytes.decode("latin-1")
        except Exception as e:
            raise ValueError(f"Error decoding text notes: {str(e)}")
    return clean_text(text_content)

# Trigger reload




